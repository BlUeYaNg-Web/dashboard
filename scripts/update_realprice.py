#!/usr/bin/env python3
"""
update_realprice.py  v2  2026-05-31
每週日 13:00 本機執行

資料來源：Google Sheets（新版 1obKJ7DT__0MeclJGJc7cNzQ8eYLz6T2kX5JdcrVsINs）
欄位：登錄日期, 買賣日期, 所在區域, 建案名稱, 戶別樓層, 平均單價, ...坪數...

修正：
  - 新 Sheet floor/unit 格式解析（支援「B/7/二樓」→B7、「A/A5-05F/五樓」→A5 等）
  - 中文樓層轉換（二樓→2F、十一樓→11F）
  - 基準從 GitHub 讀（不依賴本機 JSON 寫入權限）
  - 推送 pptx_update.json 改用 GitHub API（不用 git push）
"""

import csv, json, os, re, glob, sys, base64, urllib.request
import requests
from datetime import datetime, timedelta
from io import StringIO
from pptx import Presentation
from pptx.dml.color import RGBColor

sys.stdout.reconfigure(encoding='utf-8', errors='replace')

# ───────── 設定 ─────────────────────────────────────────────────────────────
SHEET_ID  = "1obKJ7DT__0MeclJGJc7cNzQ8eYLz6T2kX5JdcrVsINs"
SHEET_URL = f"https://docs.google.com/spreadsheets/d/{SHEET_ID}/export?format=csv"

PPTX_DIR  = r"I:\我的雲端硬碟\家泰嘉潤丰藏JANU\03.建設公司\01.業主回報\週報"
PPTX_BASE = "土城區實價登錄明細"

cfg   = json.loads(open(r'C:\Users\dxb71\.github_dashboard.json').read())
TOKEN = cfg['token']
REPO  = cfg['repo']

BUILDINGS = [
    ("丰藏",    "丰藏"),
    ("台信琢岳", "台信琢岳"),
    ("紅布朗",  "紅布朗花園"),
    ("寶佳淳青", "寶佳淳青"),
    ("松陽馥麗", "松陽馥麗"),
    ("新濠岳",  "新濠岳"),
    ("朗沐",    "朗沐"),
    ("合康雙匯", "合康雙匯"),
    ("和耀美家", "和耀美家雅居"),
    ("迴東騰",  "迴東騰"),
    ("天好運",  "天好運"),
    ("若水秧翠", "若水秧翠"),
]
MONTH_CN = {
    "一月":1,"二月":2,"三月":3,"四月":4,"五月":5,"六月":6,
    "七月":7,"八月":8,"九月":9,"十月":10,"十一月":11,"十二月":12,
}
CN_TO_INT = {
    '一':1,'二':2,'三':3,'四':4,'五':5,'六':6,'七':7,'八':8,'九':9,
    '十':10,'十一':11,'十二':12,'十三':13,'十四':14,'十五':15,
    '十六':16,'十七':17,'十八':18,'十九':19,'二十':20,
    '二十一':21,'二十二':22,'二十三':23,'二十四':24,'二十五':25,
}
RED   = RGBColor(0xFF, 0x00, 0x00)
BLACK = RGBColor(0x00, 0x00, 0x00)


# ───────── 工具函式 ──────────────────────────────────────────────────────────
def roc_to_date(s):
    try:
        s = str(int(float(str(s).strip())))
        if len(s) == 7:
            return datetime(int(s[:3]) + 1911, int(s[3:5]), int(s[5:7]))
    except Exception:
        pass
    return None


def parse_price(s):
    m = re.match(r'([\d.]+)萬', str(s).strip())
    return round(float(m.group(1)), 2) if m else None


def chinese_floor_to_nf(floor_str):
    """'二樓' → '2F', '十一樓' → '11F'"""
    s = floor_str.replace('樓', '').strip()
    n = CN_TO_INT.get(s)
    if n:
        return f"{n}F"
    # fallback：數字+F
    m = re.match(r'(\d+)[Ff]?', s)
    if m:
        return f"{int(m.group(1))}F"
    return None


def parse_floor_unit(floor_unit_str, col_headers):
    """
    解析 Google Sheet 的「戶別樓層」欄位，回傳 (floor_label, unit_label)。
    floor_label 如 '2F', unit_label 如 'H' / 'A2' / 'B7'。

    觀察到的格式：
      "H/02F/二樓"    → unit=H,  floor=2F
      "A/2/七樓"      → unit=A2, floor=7F  (A+2→A2)
      "C2/03F/三樓"   → unit=C2, floor=3F
      "B/7/二樓"      → unit=B7, floor=2F  (B+7→B7)
      "/A7-4F/四樓"   → unit=A7, floor=4F
      "A/A5-05F/五樓" → unit=A5, floor=5F
      "A3/8F/八樓"    → unit=A3, floor=8F
      "A2/11/十一樓"  → unit=A2, floor=11F
    """
    parts = [p.strip() for p in floor_unit_str.split('/')]
    if len(parts) < 3:
        return None, None

    p0, p1, p2 = parts[0], parts[1], parts[2]

    # ── 樓層：優先從中文（p2），fallback 取 p1 的數字
    floor_label = chinese_floor_to_nf(p2)
    if not floor_label:
        fm = re.search(r'(\d+)[Ff]', p1)
        if fm:
            floor_label = f"{int(fm.group(1))}F"
        else:
            return None, None

    # ── 戶別：依序嘗試候選，找到符合 PPTX 欄位的即採用
    col_up = [h.upper() for h in col_headers]

    def try_unit(cand):
        if cand and cand.upper() in col_up:
            return col_headers[col_up.index(cand.upper())]
        return None

    candidates = []
    if p0:
        candidates.append(p0)  # "H", "C2", "A3", "A2" 直接匹配
        # 嘗試 p0 + p1 的純數字部分（"B"+"7"="B7", "A"+"2"="A2"）
        num_m = re.match(r'^(\d+)', p1)
        if num_m:
            candidates.append(p0 + num_m.group(1))

    # 從 p1 提取字母+數字前綴（"A7-4F"→"A7", "A5-05F"→"A5", "02F"→忽略）
    alpha_m = re.match(r'([A-Za-z]+\d*)', p1)
    if alpha_m and not p1[0].isdigit():
        candidates.append(alpha_m.group(1))

    for cand in candidates:
        unit = try_unit(cand)
        if unit:
            return floor_label, unit

    return floor_label, None


# ───────── 資料抓取 ──────────────────────────────────────────────────────────
def fetch_sheets():
    r = requests.get(SHEET_URL, headers={'User-Agent': 'Mozilla/5.0'}, timeout=30)
    r.raise_for_status()
    r.encoding = 'utf-8-sig'
    rows = list(csv.reader(StringIO(r.text)))
    return rows[2:] if len(rows) > 2 else []   # 跳過 2 列 header


def get_baseline_from_github():
    """從 GitHub 讀取上週基準（各建案 total_records）"""
    url = f'https://api.github.com/repos/{REPO}/contents/results/pptx_update.json'
    req = urllib.request.Request(url, headers={'Authorization': f'token {TOKEN}'})
    with urllib.request.urlopen(req) as resp:
        d = json.load(resp)
    prev = json.loads(base64.b64decode(d['content']).decode('utf-8'))
    print(f"  基準版本：{prev.get('run_at','?')}")
    return {c['name']: c['total_records'] for c in prev.get('cases', [])}, d['sha']


def push_json_to_github(data, sha):
    url = f'https://api.github.com/repos/{REPO}/contents/results/pptx_update.json'
    b64 = base64.b64encode(json.dumps(data, ensure_ascii=False, indent=2).encode()).decode()
    payload = json.dumps({'message': f'Update 實價登錄 {data["run_at"]}',
                          'content': b64, 'sha': sha}).encode()
    req = urllib.request.Request(url, data=payload, method='PUT',
          headers={'Authorization': f'token {TOKEN}', 'Content-Type': 'application/json'})
    with urllib.request.urlopen(req) as resp:
        return json.load(resp)['commit']['sha'][:8]


# ───────── 建案統計 ──────────────────────────────────────────────────────────
def get_building_stats(rows, keyword, prev_total):
    monthly, all_txs = {}, []
    for row in rows:
        if len(row) < 6:
            continue
        if keyword not in row[3]:
            continue
        reg_date = roc_to_date(row[0])
        if not reg_date:
            continue
        price      = parse_price(row[5])
        floor_unit = row[4].strip() if len(row) > 4 else ""
        area_raw   = row[8].strip() if len(row) > 8 else ""
        area_m     = re.match(r'([\d.]+)', area_raw)
        area       = round(float(area_m.group(1)), 2) if area_m else None
        monthly[reg_date.month] = monthly.get(reg_date.month, 0) + 1
        all_txs.append({'date': reg_date, 'price': price,
                        'floor_unit': floor_unit, 'area': area})

    all_txs.sort(key=lambda x: x['date'])
    current_total = len(all_txs)
    weekly_new    = max(0, current_total - prev_total) if prev_total is not None else 0
    weekly_txs    = all_txs[-weekly_new:] if weekly_new > 0 else []
    latest        = all_txs[-1] if all_txs else {}

    if all_txs:
        weeks_dry = 0 if weekly_new > 0 else max(0, (datetime.now() - all_txs[-1]['date']).days // 7)
    else:
        weeks_dry = None

    return {
        'weekly_new': weekly_new, 'current_total': current_total,
        'monthly': monthly, 'weekly_txs': weekly_txs,
        'latest_price': latest.get('price'), 'latest_area': latest.get('area'),
        'weeks_dry': weeks_dry,
    }


# ───────── PPTX 更新 ─────────────────────────────────────────────────────────
def update_weekly_text(slide, weekly_new):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        ft = shape.text_frame.text
        if '本週新增' not in ft or '戶' not in ft:
            continue
        for para in shape.text_frame.paragraphs:
            runs = para.runs
            for i, run in enumerate(runs):
                if '本週新增' in run.text and i + 1 < len(runs):
                    nr = runs[i + 1]
                    nr.text = str(weekly_new)
                    nr.font.color.rgb = RED if weekly_new > 0 else BLACK
                    break
            else:
                for run in runs:
                    if '本週新增' in run.text and '戶' in run.text:
                        run.text = re.sub(r'本週新增\s*\d+\s*戶',
                                          f'本週新增 {weekly_new} 戶', run.text)
                        run.font.color.rgb = RED if weekly_new > 0 else BLACK


def update_monthly_shapes(slide, monthly):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        for cn, num in MONTH_CN.items():
            if text.startswith(cn):
                count = monthly.get(num, 0)
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = re.sub(rf'({cn}[：:]\s*)\d+', rf'\g<1>{count}', run.text)
                break


def _get_font_tmpl(table):
    from pptx.enum.text import PP_ALIGN
    fn, fs, bold, align = '微軟正黑體', None, None, PP_ALIGN.CENTER
    for r, row in enumerate(table.rows):
        if r < 2:
            continue
        for c, cell in enumerate(row.cells):
            if c == 0 or '萬' not in cell.text:
                continue
            tf = cell.text_frame
            if tf.paragraphs:
                align = tf.paragraphs[0].alignment or align
                for run in tf.paragraphs[0].runs:
                    return run.font.name or fn, run.font.size, run.font.bold, align
    return fn, fs, bold, align


def _write_cell(cell, text, fn, fs, bold, align, color=None):
    tf = cell.text_frame
    tf.word_wrap = False
    for para in tf.paragraphs[1:]:
        para._p.getparent().remove(para._p)
    para = tf.paragraphs[0]
    para.alignment = align
    for run in para.runs:
        run._r.getparent().remove(run._r)
    run = para.add_run()
    run.text = text
    run.font.name = fn
    if fs is not None:
        run.font.size = fs
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def update_price_table(slide, weekly_txs):
    """把本週新增的成交格子標紅色字體"""
    total_updated = 0
    for shape in slide.shapes:
        if shape.shape_type != 19:
            continue
        if shape.name in ('表格 10', '表格 7'):
            continue
        table = shape.table
        col_hdrs = [c.text.strip() for c in table.rows[0].cells]
        row_hdrs = [table.rows[r].cells[0].text.strip() for r in range(len(table.rows))]
        fn, fs, bold, align = _get_font_tmpl(table)

        for tx in weekly_txs:
            fu = tx.get('floor_unit', '')
            price = tx.get('price')
            if not price:
                continue

            floor_label, unit_label = parse_floor_unit(fu, col_hdrs)
            if not floor_label or not unit_label:
                print(f"    ⚠ 無法解析：{fu!r}  (floor={floor_label}, unit={unit_label})")
                continue

            ri = next((i for i, h in enumerate(row_hdrs) if h == floor_label), None)
            ci = next((i for i, h in enumerate(col_hdrs) if h == unit_label), None)

            if ri is not None and ci is not None:
                _write_cell(table.rows[ri].cells[ci],
                            f'{price}萬', fn, fs, bold, align, RED)
                print(f"    ★ 標紅 [{floor_label}][{unit_label}] → {price}萬")
                total_updated += 1
            else:
                print(f"    ✗ 找不到格子 floor={floor_label} unit={unit_label}  (fu={fu!r})")

    return total_updated


# ───────── 主流程 ────────────────────────────────────────────────────────────
def main():
    today    = datetime.now()
    run_at   = today.strftime('%Y/%m/%d %H:%M')
    date_tag = today.strftime('%y%m%d')

    print(f"[{run_at}] 開始更新實價登錄（v2）")

    # Step 1: 基準
    print("\nStep 1: 從 GitHub 讀取上週基準...")
    prev_totals, json_sha = get_baseline_from_github()

    # Step 2: Google Sheet
    print("\nStep 2: 下載 Google Sheets...")
    rows = fetch_sheets()
    print(f"  取得 {len(rows)} 筆")

    # Step 3: PPTX 底稿（用最新手動版）
    all_pptx = sorted(glob.glob(os.path.join(PPTX_DIR, f'{PPTX_BASE}*.pptx')))
    manual   = [f for f in all_pptx if '自動更新' not in f]
    src_path = manual[-1] if manual else all_pptx[-1]
    print(f"\nStep 3: 底稿 → {os.path.basename(src_path)}")
    prs = Presentation(src_path)

    # Step 4: 逐建案計算
    print("\nStep 4: 計算本週新增...")
    cases, total_red = [], 0
    for keyword, display in BUILDINGS:
        prev = prev_totals.get(display)
        stats = get_building_stats(rows, keyword, prev)
        wn = stats['weekly_new']
        print(f"\n  [{display}] 基準={prev} → 現在={stats['current_total']} → 新增={wn}戶")
        if wn > 0:
            for tx in stats['weekly_txs']:
                print(f"    {tx['date'].strftime('%Y/%m/%d')} {tx['floor_unit']} {tx['price']}萬")

        for slide in prs.slides:
            slide_text = ' '.join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
            if display in slide_text or keyword in slide_text:
                update_weekly_text(slide, wn)
                update_monthly_shapes(slide, stats['monthly'])
                if stats['weekly_txs']:
                    total_red += update_price_table(slide, stats['weekly_txs'])

        cases.append({
            'name': display,
            'weekly_new': wn,
            'total_records': stats['current_total'],
            'latest_price': stats['latest_price'],
            'latest_area': stats['latest_area'],
            'weeks_dry': stats['weeks_dry'],
        })

    # Step 5: 存 PPTX
    new_name = f'{PPTX_BASE}{date_tag}自動更新.pptx'
    new_path = os.path.join(PPTX_DIR, new_name)
    prs.save(new_path)
    print(f"\nStep 5: PPTX 已存 → {new_name}（標紅 {total_red} 格）")

    # Step 6: 推送 JSON
    json_data = {'run_at': run_at, 'output_file': new_path, 'cases': cases}
    commit = push_json_to_github(json_data, json_sha)
    print(f"\nStep 6: pptx_update.json 推送 GitHub (commit: {commit})")

    print(f"\n=== 完成 ===  標紅 {total_red} 格，本週新增 {sum(c['weekly_new'] for c in cases)} 戶")
    for c in cases:
        if c['weekly_new'] > 0:
            print(f"  ★ {c['name']}: +{c['weekly_new']} 戶")


if __name__ == '__main__':
    main()
